"""Validate build_from_yaml(): same content.yaml → correct .pptx."""
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Emu

from build_ppt import build_from_yaml

SKILL_DIR = Path(__file__).resolve().parent.parent
SAMPLE = SKILL_DIR / "sample" / "content.yaml"


@pytest.fixture(scope="module")
def built(tmp_path_factory):
    if not SAMPLE.exists():
        pytest.skip("sample/content.yaml missing")
    out = tmp_path_factory.mktemp("pptx") / "out.pptx"
    build_from_yaml(str(SAMPLE), str(out))
    return out


def test_slide_count_matches_yaml(built):
    import yaml
    expected = len(yaml.safe_load(SAMPLE.read_text())["slides"])
    p = Presentation(str(built))
    assert len(p.slides) == expected


def test_template_chrome_preserved(built):
    """Red banner shape from slideMaster must remain unchanged."""
    p = Presentation(str(built))
    # The banner lives in slideMaster, not the slides themselves.
    # Pull it back via the master shapes to verify it survived.
    master = p.slide_masters[0]
    banners = [
        s for s in master.shapes
        if s.shape_type == 1 and s.height == Emu(914400)  # 1 inch = banner
    ]
    assert len(banners) >= 1, "red banner shape missing from slideMaster"


def test_first_slide_title_matches_yaml(built):
    """First slide should contain the YAML title text somewhere in its shapes.

    The template layout puts subtitle and title in two separate placeholders
    whose XML order is implementation-defined, so we scan all text shapes
    rather than picking the first.
    """
    import yaml
    spec = yaml.safe_load(SAMPLE.read_text())
    p = Presentation(str(built))
    first_slide = p.slides[0]
    all_text = " ".join(
        shape.text_frame.text
        for shape in first_slide.shapes
        if shape.has_text_frame and shape.text_frame.text
    )
    assert spec["slides"][0]["title"] in all_text
