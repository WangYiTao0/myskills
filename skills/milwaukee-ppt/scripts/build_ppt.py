"""Milwaukee Tool brand deck builder (python-pptx wrapper).

Usage:
    from build_ppt import MilwaukeeDeck

    deck = MilwaukeeDeck()
    s1 = deck.add_slide("PRODUCT OVERVIEW", "M18 FUEL Series")
    s1.add_paragraphs([
        ("产品定位", {"bold": True, "size": 18}),
        ("面向专业用户的高扭矩冲击钻系列。", {"size": 14}),
    ])

    s2 = deck.add_slide("KEY SPECS", "Performance")
    s2.add_table(
        [["Model", "Torque", "RPM"],
         ["M18-1", "1000 Nm", "0-2000"],
         ["M18-2", "1200 Nm", "0-2100"]],
        header=True,
    )
    s2.add_image("/path/to/product.png", left_cm=2, top_cm=10, width_cm=10)

    deck.save("output.pptx")

Coordinates are EMU under the hood. Helpers accept cm via *_cm kwargs.
1 cm = 360000 EMU. Slide is 33.87 cm x 19.05 cm (16:9).
"""
from __future__ import annotations

from pathlib import Path
from typing import Iterable

from pptx import Presentation
from pptx.util import Cm, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree


SKILL_DIR = Path(__file__).resolve().parent.parent
DEFAULT_TEMPLATE = SKILL_DIR / "assets" / "template.pptx"

# Brand palette — see references/design-guidelines.md for usage rules.
MILWAUKEE_RED = RGBColor(0xDB, 0x02, 0x1D)
TEXT_DARK = RGBColor(0x33, 0x33, 0x33)
TEXT_MID = RGBColor(0x66, 0x66, 0x66)
BG_LIGHT = RGBColor(0xF5, 0xF5, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
# Status colors (charts / KPIs only, not for layout decoration)
STATUS_OK = RGBColor(0x2E, 0x7D, 0x32)
STATUS_WARN = RGBColor(0xED, 0x6C, 0x02)
STATUS_DANGER = RGBColor(0xC6, 0x28, 0x28)

PALETTE = {
    "red": MILWAUKEE_RED,
    "dark": TEXT_DARK,
    "mid": TEXT_MID,
    "bg_light": BG_LIGHT,
    "white": WHITE,
    "ok": STATUS_OK,
    "warn": STATUS_WARN,
    "danger": STATUS_DANGER,
}

# Font fallback chain. python-pptx font.name sets latin typeface only;
# for east-asian (zh) we set <a:ea> via direct rPr edit.
LATIN_FONT = "Calibri"
EA_FONT_PRIMARY = "Microsoft JhengHei"  # zh-TW preferred
EA_FONT_FALLBACKS = ["PingFang TC", "Microsoft YaHei", "Heiti TC", "SimHei"]

# Content area bounds (inside the white area, below red banner, above footer)
CONTENT_LEFT_CM = 1.0
CONTENT_TOP_CM = 2.5
CONTENT_RIGHT_CM = 32.87  # = slide_width 33.87 - 1.0
CONTENT_BOTTOM_CM = 17.5  # leave room for footer
CONTENT_WIDTH_CM = CONTENT_RIGHT_CM - CONTENT_LEFT_CM
CONTENT_HEIGHT_CM = CONTENT_BOTTOM_CM - CONTENT_TOP_CM


def _set_run_font(run, *, name=None, size=None, bold=None, color=None, ea=None):
    """Set run formatting including east-asian typeface (which python-pptx
    doesn't expose directly)."""
    if name is not None:
        run.font.name = name
    if size is not None:
        run.font.size = Pt(size)
    if bold is not None:
        run.font.bold = bold
    if color is not None:
        run.font.color.rgb = color
    if ea is not None:
        rPr = run._r.get_or_add_rPr()
        for tag in ("a:ea",):
            existing = rPr.find(qn(tag))
            if existing is not None:
                rPr.remove(existing)
        ea_el = etree.SubElement(rPr, qn("a:ea"))
        ea_el.set("typeface", ea)


class _Slide:
    def __init__(self, slide):
        self._slide = slide
        self._next_top_cm = CONTENT_TOP_CM

    # ---- placeholders ----
    def set_title(self, text: str):
        ph = self._slide.placeholders[0]
        ph.text_frame.text = text
        return self

    def set_subtitle(self, text: str):
        ph = self._slide.placeholders[10]
        ph.text_frame.text = text
        return self

    # ---- content helpers ----
    def add_paragraphs(
        self,
        items: Iterable[tuple[str, dict]],
        *,
        left_cm: float | None = None,
        top_cm: float | None = None,
        width_cm: float | None = None,
        height_cm: float | None = None,
    ):
        """Add a textbox of paragraphs.
        items: iterable of (text, style_dict). style_dict keys:
          size (pt), bold, color (RGBColor), align ('l'|'c'|'r'),
          bullet (bool), ea (str, east-asian typeface).
        """
        left = Cm(left_cm if left_cm is not None else CONTENT_LEFT_CM)
        top = Cm(top_cm if top_cm is not None else self._next_top_cm)
        width = Cm(width_cm if width_cm is not None else CONTENT_WIDTH_CM)
        height = Cm(height_cm if height_cm is not None else 5.0)

        tb = self._slide.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame
        tf.word_wrap = True

        for i, (text, style) in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            align = style.get("align")
            if align:
                p.alignment = {"l": PP_ALIGN.LEFT, "c": PP_ALIGN.CENTER, "r": PP_ALIGN.RIGHT}[align]
            if style.get("bullet"):
                _apply_bullet(p)
            run = p.add_run()
            run.text = text
            _set_run_font(
                run,
                name=style.get("name", LATIN_FONT),
                size=style.get("size", 14),
                bold=style.get("bold"),
                color=style.get("color", TEXT_DARK),
                ea=style.get("ea", EA_FONT_PRIMARY),
            )

        self._next_top_cm = (top + height) / 360000 + 0.3
        return tb

    def add_bullets(
        self,
        bullets: list[str],
        *,
        size: int = 14,
        **box_kwargs,
    ):
        return self.add_paragraphs(
            [(b, {"size": size, "bullet": True}) for b in bullets],
            **box_kwargs,
        )

    def add_table(
        self,
        rows: list[list[str]],
        *,
        header: bool = True,
        left_cm: float | None = None,
        top_cm: float | None = None,
        width_cm: float | None = None,
        height_cm: float | None = None,
    ):
        n_rows = len(rows)
        n_cols = max(len(r) for r in rows)
        left = Cm(left_cm if left_cm is not None else CONTENT_LEFT_CM)
        top = Cm(top_cm if top_cm is not None else self._next_top_cm)
        width = Cm(width_cm if width_cm is not None else CONTENT_WIDTH_CM)
        height = Cm(height_cm if height_cm is not None else 0.8 * n_rows + 0.4)

        gf = self._slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
        table = gf.table
        for r, row in enumerate(rows):
            for c in range(n_cols):
                cell = table.cell(r, c)
                cell.text = row[c] if c < len(row) else ""
                p = cell.text_frame.paragraphs[0]
                run = p.runs[0]
                if header and r == 0:
                    cell.fill.solid()
                    cell.fill.fore_color.rgb = MILWAUKEE_RED
                    _set_run_font(run, name=LATIN_FONT, size=12, bold=True, color=WHITE, ea=EA_FONT_PRIMARY)
                else:
                    _set_run_font(run, name=LATIN_FONT, size=11, color=TEXT_DARK, ea=EA_FONT_PRIMARY)

        self._next_top_cm = (top + height) / 360000 + 0.3
        return gf

    def add_image(
        self,
        path: str | Path,
        *,
        left_cm: float | None = None,
        top_cm: float | None = None,
        width_cm: float | None = None,
        height_cm: float | None = None,
    ):
        kwargs = {}
        if width_cm is not None:
            kwargs["width"] = Cm(width_cm)
        if height_cm is not None:
            kwargs["height"] = Cm(height_cm)
        left = Cm(left_cm if left_cm is not None else CONTENT_LEFT_CM)
        top = Cm(top_cm if top_cm is not None else self._next_top_cm)
        pic = self._slide.shapes.add_picture(str(path), left, top, **kwargs)
        self._next_top_cm = (pic.top + pic.height) / 360000 + 0.3
        return pic

    # ---- polish helpers ----
    def columns(self, n: int, *, ratios: list[float] | None = None,
                gap_cm: float = 0.5,
                top_cm: float | None = None,
                height_cm: float | None = None) -> list[dict]:
        """Return n column rects covering the content area.

        Each rect is a dict {left_cm, top_cm, width_cm, height_cm} ready to be
        passed as **kwargs to add_paragraphs/add_image/add_table.

        Example:
            cols = slide.columns(3)
            for col, item in zip(cols, items):
                slide.add_paragraphs([(item, {"size": 16})], **col)
        """
        if n < 1:
            raise ValueError("n must be >= 1")
        ratios = ratios or [1.0] * n
        if len(ratios) != n:
            raise ValueError("len(ratios) must equal n")
        total_gap = gap_cm * (n - 1)
        usable = CONTENT_WIDTH_CM - total_gap
        unit = usable / sum(ratios)
        widths = [r * unit for r in ratios]
        top = top_cm if top_cm is not None else self._next_top_cm
        height = height_cm if height_cm is not None else (CONTENT_BOTTOM_CM - top)
        rects = []
        x = CONTENT_LEFT_CM
        for w in widths:
            rects.append({"left_cm": x, "top_cm": top, "width_cm": w, "height_cm": height})
            x += w + gap_cm
        return rects

    def add_kpi(self, value: str, label: str, *,
                color: RGBColor = MILWAUKEE_RED,
                left_cm: float | None = None,
                top_cm: float | None = None,
                width_cm: float | None = None):
        """Big number + small caption underneath. Use for stat highlights."""
        left = Cm(left_cm if left_cm is not None else CONTENT_LEFT_CM)
        top = Cm(top_cm if top_cm is not None else self._next_top_cm)
        width = Cm(width_cm if width_cm is not None else CONTENT_WIDTH_CM)
        height = Cm(4.5)
        tb = self._slide.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame
        tf.word_wrap = True
        p1 = tf.paragraphs[0]
        p1.alignment = PP_ALIGN.CENTER
        r1 = p1.add_run()
        r1.text = value
        _set_run_font(r1, name=LATIN_FONT, size=72, bold=True, color=color, ea=EA_FONT_PRIMARY)
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = label
        _set_run_font(r2, name=LATIN_FONT, size=14, color=TEXT_MID, ea=EA_FONT_PRIMARY)
        self._next_top_cm = (top + height) / 360000 + 0.3
        return tb

    def add_quote(self, text: str, author: str = "", *,
                  left_cm: float | None = None,
                  top_cm: float | None = None,
                  width_cm: float | None = None):
        """Centered italic quote with optional right-aligned attribution."""
        left = Cm(left_cm if left_cm is not None else CONTENT_LEFT_CM + 2.0)
        top = Cm(top_cm if top_cm is not None else self._next_top_cm)
        width = Cm(width_cm if width_cm is not None else CONTENT_WIDTH_CM - 4.0)
        height = Cm(5.0)
        tb = self._slide.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = f"“{text}”"
        _set_run_font(r, name=LATIN_FONT, size=24, color=TEXT_DARK, ea=EA_FONT_PRIMARY)
        r.font.italic = True
        if author:
            p2 = tf.add_paragraph()
            p2.alignment = PP_ALIGN.RIGHT
            r2 = p2.add_run()
            r2.text = f"— {author}"
            _set_run_font(r2, name=LATIN_FONT, size=14, color=TEXT_MID, ea=EA_FONT_PRIMARY)
        self._next_top_cm = (top + height) / 360000 + 0.3
        return tb

    def add_section_divider(self, text: str):
        """Big centered red text — for chapter intros. Use as a full-page block."""
        left = Cm(CONTENT_LEFT_CM)
        top = Cm(CONTENT_TOP_CM + 4.0)
        width = Cm(CONTENT_WIDTH_CM)
        height = Cm(6.0)
        tb = self._slide.shapes.add_textbox(left, top, width, height)
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = text
        _set_run_font(r, name=LATIN_FONT, size=54, bold=True, color=MILWAUKEE_RED, ea=EA_FONT_PRIMARY)
        self._next_top_cm = (top + height) / 360000 + 0.3
        return tb

    def add_speaker_notes(self, text: str):
        """Add speaker notes (operational text that should not appear on slide face)."""
        notes = self._slide.notes_slide
        notes.notes_text_frame.text = text
        return notes


def _apply_bullet(paragraph):
    """Add a • bullet to a paragraph via direct pPr edit (python-pptx
    doesn't expose buChar)."""
    pPr = paragraph._pPr
    if pPr is None:
        pPr = paragraph._p.get_or_add_pPr()
    pPr.set("marL", "342900")
    pPr.set("indent", "-342900")
    for tag in ("a:buChar", "a:buAutoNum", "a:buNone"):
        existing = pPr.find(qn(tag))
        if existing is not None:
            pPr.remove(existing)
    bu = etree.SubElement(pPr, qn("a:buChar"))
    bu.set("char", "•")


class MilwaukeeDeck:
    def __init__(self, template_path: str | Path = DEFAULT_TEMPLATE):
        template_path = Path(template_path)
        if not template_path.exists():
            raise FileNotFoundError(f"template not found: {template_path}")
        self._prs = Presentation(str(template_path))
        # Reuse the bundled slide1 as the first page; subsequent pages added
        # via the blank layout (which carries the same title/body placeholders).
        self._first_used = False
        self._layout = self._prs.slide_layouts[0]

    def add_slide(self, title: str, subtitle: str = "") -> _Slide:
        if not self._first_used:
            slide = self._prs.slides[0]
            self._first_used = True
        else:
            slide = self._prs.slides.add_slide(self._layout)
        wrapper = _Slide(slide)
        wrapper.set_title(title)
        if subtitle:
            wrapper.set_subtitle(subtitle)
        return wrapper

    def save(self, path: str | Path, *, force: bool = False) -> Path:
        """Save the deck. Refuses to overwrite a file that PowerPoint has open
        (detected via the platform-standard ~$ lock file). Pass force=True to
        skip the check."""
        path = Path(path).resolve()
        path.parent.mkdir(parents=True, exist_ok=True)
        if not force:
            lock = path.parent / f"~${path.name}"
            if lock.exists():
                raise RuntimeError(
                    f"{path.name} appears to be open in PowerPoint (lock file "
                    f"{lock.name} exists). Close it first, or pass force=True."
                )
        self._prs.save(str(path))
        return path


# ---------------------------------------------------------------------------
# YAML entry point (added 2026-05-18 — dual-track refactor)
# ---------------------------------------------------------------------------

def _promote_title_placeholder(slide) -> None:
    """Move the title placeholder element (ph type='title') to be the first
    sp child of spTree so that iterating slide.shapes yields the title before
    the subtitle.  The template stores subtitle (ph idx=10) before title in
    the XML; this fixes the ordering without touching the template binary."""
    spTree = slide._element.spTree
    from pptx.oxml.ns import qn as _qn
    title_sp = None
    for sp in list(spTree):
        ph_el = sp.find(_qn("p:ph"))
        if ph_el is None:
            # Try nested inside nvSpPr/nvPr
            ph_el = sp.find(f".//{_qn('p:ph')}")
        if ph_el is not None:
            # Title placeholder: type="title" OR (no idx / idx="0")
            ph_type = ph_el.get("type", "")
            ph_idx = ph_el.get("idx")
            if ph_type == "title" or (ph_idx is None or ph_idx == "0"):
                title_sp = sp
                break
    if title_sp is not None:
        spTree.remove(title_sp)
        # nvGrpSpPr is [0], grpSpPr is [1]; insert title as first sp at [2]
        spTree.insert(2, title_sp)


def _render_slide_from_spec(deck: "MilwaukeeDeck", spec: dict) -> None:
    """Dispatch one slide spec to the appropriate MilwaukeeDeck API."""
    t = spec["type"]
    title = spec.get("title", "")
    subtitle = spec.get("subtitle", "")
    s = deck.add_slide(title, subtitle)
    _promote_title_placeholder(s._slide)

    if t == "title":
        return

    if t == "text":
        items = []
        for block in spec.get("blocks", []):
            style = {k: v for k, v in block.items() if k != "text"}
            items.append((block["text"], style))
        s.add_paragraphs(items)
    elif t == "bullets":
        s.add_bullets(spec["items"])
    elif t == "table":
        s.add_table(spec["rows"])
    elif t == "kpi":
        s.add_kpi(spec["value"], spec["label"])
    elif t == "quote":
        s.add_quote(spec["text"], spec.get("author", ""))
    elif t == "section":
        s.add_section_divider(spec.get("text", title))
    elif t == "image":
        s.add_image(spec["path"])
        if caption := spec.get("caption"):
            s.add_paragraphs([(caption, {"size": 14, "color": TEXT_MID})])
    elif t == "columns":
        cols = spec["columns"]
        rects = s.columns(len(cols))
        for col, rect in zip(cols, rects):
            s.add_paragraphs(
                [
                    (col["head"], {"size": 20, "bold": True, "color": MILWAUKEE_RED}),
                    (col["body"], {"size": 16}),
                ],
                **rect,
            )
    else:
        raise ValueError(f"unknown slide type {t!r}")


def build_from_yaml(content_path: str, out_path: str,
                    template_path: str | None = None) -> str:
    """YAML → .pptx one-shot entry. Validates schema, dispatches each slide
    to the matching MilwaukeeDeck API, saves the result.

    Returns the absolute output path.
    """
    import yaml
    import sys
    sys.path.insert(0, str(Path(__file__).resolve().parent))
    from fill_html import validate_content

    content = yaml.safe_load(Path(content_path).read_text())
    validate_content(content)

    deck = MilwaukeeDeck(template_path or DEFAULT_TEMPLATE)
    for spec in content["slides"]:
        _render_slide_from_spec(deck, spec)
    return deck.save(out_path)


def _cli() -> int:
    import argparse
    p = argparse.ArgumentParser(description="Milwaukee YAML → .pptx")
    p.add_argument("content", help="path to content.yaml")
    p.add_argument("--out", required=True, help="output .pptx path")
    p.add_argument("--template", default=None, help="override template.pptx")
    args = p.parse_args()
    out = build_from_yaml(args.content, args.out, args.template)
    print(f"wrote {out}")
    return 0


if __name__ == "__main__":
    import sys
    sys.exit(_cli())
