"""Polish linter for Milwaukee Tool decks.

Reads a built .pptx and reports violations of the design guidelines
(see references/design-guidelines.md). Designed to be run as the final
QA step before delivery.

Usage:
    python polish.py output.pptx
    python polish.py output.pptx --strict   # exit non-zero on any warning
"""
from __future__ import annotations

import argparse
import re
import sys
from collections import Counter
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


# Acceptable colors expressed as 0xRRGGBB ints.
PALETTE_INTS = {
    0xDB021D,  # red
    0x333333, 0x666666, 0x999999,  # dark/mid/light gray
    0xF5F5F5, 0xFFFFFF, 0x000000,  # backgrounds
    0x2E7D32, 0xED6C02, 0xC62828,  # status
}

ALLOWED_FONT_FAMILIES = {
    "Calibri",
    "Microsoft JhengHei", "Microsoft YaHei",
    "PingFang TC", "PingFang SC",
    "Heiti TC", "Heiti SC", "SimHei",
    "+mn-lt", "+mj-lt", "+mn-ea", "+mj-ea",  # theme references
}

# Banner / footer cutoff — anything wholly inside these zones is template chrome.
BANNER_BOTTOM_EMU = 900_000   # 2.5 cm; placeholders are above this
FOOTER_TOP_EMU = 6_500_000    # 18.05 cm

MAX_TITLE_CHARS = 40
MAX_SUBTITLE_CHARS = 60
MAX_BULLETS_PER_SLIDE = 6
MAX_WORDS_PER_BULLET_EN = 6
MAX_CHARS_PER_BULLET_ZH = 12
MAX_TABLE_ROWS = 7
MAX_FONT_FAMILIES = 2
MIN_IMAGE_PIXELS = 1280 * 720


class Report:
    def __init__(self):
        self.errors: list[str] = []
        self.warnings: list[str] = []
        self.info: list[str] = []

    def err(self, slide_no, msg):
        self.errors.append(f"slide {slide_no}: {msg}")

    def warn(self, slide_no, msg):
        self.warnings.append(f"slide {slide_no}: {msg}")

    def note(self, msg):
        self.info.append(msg)

    def print(self):
        for m in self.info:
            print(f"  ℹ  {m}")
        for m in self.warnings:
            print(f"  ⚠  {m}")
        for m in self.errors:
            print(f"  ✖  {m}")
        total = len(self.errors) + len(self.warnings)
        print(f"\n{len(self.errors)} errors, {len(self.warnings)} warnings.")
        return total


def _is_chinese(s: str) -> bool:
    return bool(re.search(r"[一-鿿]", s))


def _bullet_too_long(text: str) -> tuple[bool, str]:
    text = text.strip()
    if not text:
        return False, ""
    if _is_chinese(text):
        if len(text) > MAX_CHARS_PER_BULLET_ZH:
            return True, f"项目符号过长 ({len(text)} 字 > {MAX_CHARS_PER_BULLET_ZH})"
    else:
        words = len(text.split())
        if words > MAX_WORDS_PER_BULLET_EN:
            return True, f"bullet too long ({words} words > {MAX_WORDS_PER_BULLET_EN})"
    return False, ""


def _collect_fonts(slide) -> set[str]:
    fonts: set[str] = set()
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                name = run.font.name
                if name:
                    fonts.add(name)
    return fonts


def _collect_off_palette_colors(slide) -> set[int]:
    """Return RGB ints found in run colors that aren't in the palette."""
    off: set[int] = set()
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                try:
                    rgb = run.font.color.rgb
                except (AttributeError, TypeError, KeyError):
                    rgb = None
                if rgb is None:
                    continue
                v = int(str(rgb), 16)
                if v not in PALETTE_INTS:
                    off.add(v)
    return off


def _is_in_template_zone(shape) -> bool:
    """Shapes wholly inside the red banner or footer band."""
    if shape.top is None or shape.height is None:
        return False
    bottom = shape.top + shape.height
    if bottom <= BANNER_BOTTOM_EMU:
        return True
    if shape.top >= FOOTER_TOP_EMU:
        return True
    return False


def lint(pptx_path: Path) -> Report:
    rep = Report()
    if not pptx_path.exists():
        rep.errors.append(f"file not found: {pptx_path}")
        return rep

    lock = pptx_path.parent / f"~${pptx_path.name}"
    if lock.exists():
        rep.warn(0, f"PowerPoint lock file present ({lock.name}); file is open")

    prs = Presentation(str(pptx_path))
    rep.note(f"deck: {pptx_path.name}  slides: {len(prs.slides)}")

    for i, slide in enumerate(prs.slides, 1):
        # Title / subtitle length
        try:
            title_text = slide.placeholders[0].text_frame.text
            if len(title_text) > MAX_TITLE_CHARS:
                rep.warn(i, f"title too long ({len(title_text)} > {MAX_TITLE_CHARS}): {title_text!r}")
        except KeyError:
            pass
        try:
            sub_text = slide.placeholders[10].text_frame.text
            if len(sub_text) > MAX_SUBTITLE_CHARS:
                rep.warn(i, f"subtitle too long ({len(sub_text)} > {MAX_SUBTITLE_CHARS}): {sub_text!r}")
        except KeyError:
            pass

        # Fonts
        fonts = _collect_fonts(slide)
        unfamiliar = {f for f in fonts if f not in ALLOWED_FONT_FAMILIES}
        if unfamiliar:
            rep.warn(i, f"non-brand fonts: {sorted(unfamiliar)}")
        # Count families: collapse JhengHei/PingFang/Heiti as one CJK family,
        # Calibri as latin family.
        cjk_present = any(f for f in fonts if any(k in f for k in ("JhengHei", "PingFang", "Heiti", "YaHei", "SimHei")))
        latin_present = any(f for f in fonts if "Calibri" in f or f in {"+mn-lt", "+mj-lt"})
        family_count = sum([cjk_present, latin_present, len(unfamiliar)])
        if family_count > MAX_FONT_FAMILIES:
            rep.warn(i, f"too many font families ({family_count} > {MAX_FONT_FAMILIES})")

        # Off-palette colors
        off = _collect_off_palette_colors(slide)
        if off:
            rep.warn(i, f"off-palette colors: {[f'#{v:06X}' for v in sorted(off)]}")

        # Bullets and tables
        bullet_count = 0
        for shape in slide.shapes:
            if _is_in_template_zone(shape):
                continue
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para._pPr is not None and para._pPr.find("{http://schemas.openxmlformats.org/drawingml/2006/main}buChar") is not None:
                        bullet_count += 1
                        long_, msg = _bullet_too_long(para.text)
                        if long_:
                            rep.warn(i, msg + f": {para.text[:40]!r}")
            if shape.has_table:
                rows = len(list(shape.table.rows))
                if rows > MAX_TABLE_ROWS:
                    rep.warn(i, f"table has {rows} rows (> {MAX_TABLE_ROWS}); split or summarize")
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    img = shape.image
                    pixels = 0
                    # python-pptx exposes blob bytes; use Pillow if available
                    from io import BytesIO
                    try:
                        from PIL import Image
                        with Image.open(BytesIO(img.blob)) as im:
                            pixels = im.width * im.height
                    except ImportError:
                        pass
                    if pixels and pixels < MIN_IMAGE_PIXELS:
                        rep.warn(i, f"image below 1280x720 ({pixels} px)")
                except Exception:
                    pass

        if bullet_count > MAX_BULLETS_PER_SLIDE:
            rep.warn(i, f"too many bullets ({bullet_count} > {MAX_BULLETS_PER_SLIDE})")

    return rep


def main():
    ap = argparse.ArgumentParser()
    ap.add_argument("pptx", type=Path)
    ap.add_argument("--strict", action="store_true",
                    help="exit non-zero on any warning, not just errors")
    args = ap.parse_args()
    rep = lint(args.pptx)
    total_problems = rep.print()
    fail = len(rep.errors) > 0 or (args.strict and total_problems > 0)
    sys.exit(1 if fail else 0)


if __name__ == "__main__":
    main()
